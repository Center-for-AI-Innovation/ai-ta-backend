import asyncio
import inspect
import logging
import mimetypes
import os
import tempfile
import time
import traceback
from collections import defaultdict
from pathlib import Path
from typing import Dict, List, Union

import openai
import pytz
from dateutil import parser
from injector import inject
from langchain.embeddings.ollama import OllamaEmbeddings

# from langchain.chat_models import AzureChatOpenAI
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.schema import Document
from qdrant_client.http import models

from ai_ta_backend.database.connection_manager import ConnectionManager
from ai_ta_backend.database.sql import (
    ModelUsage,
    ProjectStats,
    WeeklyMetric,
)
from ai_ta_backend.executors.thread_pool_executor import ThreadPoolExecutorAdapter

# from ai_ta_backend.service.nomic_service import NomicService
from ai_ta_backend.service.posthog_service import PosthogService
from ai_ta_backend.service.sentry_service import SentryService

# Qwen query instruction for Illinois Chat retrieval.
# Docs are embedded without instruction during ingest; only queries get this prefix.
DEFAULT_QWEN_QUERY_INSTRUCTION = (
    "Given a user search query, retrieve the most relevant passages from the Illinois Chat knowledge "
    "base stored in the vector store to answer the query accurately. Prioritize authoritative course materials, "
    "syllabi, FAQs, official documentation, web pages, and other relevant sources. Ignore boilerplate/navigation text."
)


def _parse_allowed_embedding_providers() -> tuple[str, ...]:
    """Parse ``ALLOWED_EMBEDDING_PROVIDERS`` (comma-separated, lowercased).

    Mirrors the frontend parser in ``validation.ts``. Empty / unset → default
    ``('openai', 'ollama')`` — the two providers ``_resolve_embedding_client``
    actually knows how to build clients for.
    """
    raw = os.getenv("ALLOWED_EMBEDDING_PROVIDERS")
    if not raw:
        return ("openai", "ollama")
    parsed = tuple(p for p in (s.strip().lower() for s in raw.split(",")) if p)
    if not parsed:
        raise ValueError(
            "ALLOWED_EMBEDDING_PROVIDERS is set but parses to an empty list. "
            "Unset it or supply at least one provider."
        )
    return parsed


ALLOWED_EMBEDDING_PROVIDERS: tuple[str, ...] = _parse_allowed_embedding_providers()


class RetrievalService:
    """
    Contains all methods for business logic of the retrieval service.

    All external resource access (S3, Postgres, Qdrant/pgvector) goes through
    the injected ``ConnectionManager`` so per-project overrides are honoured.
    When a project has no per-project Qdrant override, the resolver returns a
    pgvector-default VectorDatabase whose ``upsert_main_collection`` / ``delete``
    methods transparently route to the host pgvector store.
    """

    @inject
    def __init__(
        self,
        posthog: PosthogService,
        sentry: SentryService,
        thread_pool_executor: ThreadPoolExecutorAdapter,
        conn_manager: ConnectionManager,
    ):
        self.sentry = sentry
        self.posthog = posthog
        self.thread_pool_executor = thread_pool_executor
        self.conn_manager = conn_manager

        self.openai_api_key = (
            os.getenv("OPENAI_API_KEY")
            if os.getenv("OPENAI_API_KEY")
            else os.getenv("NCSA_HOSTED_API_KEY")
        )
        self.embedding_model = (
            os.getenv("EMBEDDING_MODEL")
            if os.getenv("EMBEDDING_MODEL")
            else "text-embedding-ada-002"
        )
        self.openai_api_base = (
            os.getenv("EMBEDDING_API_BASE")
            if os.getenv("EMBEDDING_API_BASE")
            else "https://api.openai.com/v1"
        )

        if self.embedding_model == "text-embedding-ada-002":
            self.embeddings = OpenAIEmbeddings(
                model=self.embedding_model,
                openai_api_key=self.openai_api_key,
                openai_api_base=self.openai_api_base,
            )
        else:
            self.embeddings = OpenAIEmbeddings(
                model=self.embedding_model,
                openai_api_key=self.openai_api_key,
                openai_api_base=self.openai_api_base,
                # tiktoken_model_name="cl100k_base",
                tiktoken_enabled=False,
                # openai_api_key=os.environ["AZURE_OPENAI_KEY"],
                # openai_api_base=os.environ["AZURE_OPENAI_ENDPOINT"],
                # openai_api_type=os.environ['OPENAI_API_TYPE'],
                # openai_api_version=os.environ["OPENAI_API_VERSION"],
            )

        # Allow override via env; fallback to sane default for Illinois Chat retrieval.
        self.qwen_query_instruction = os.getenv(
            "QWEN_QUERY_INSTRUCTION", DEFAULT_QWEN_QUERY_INSTRUCTION
        )

        # self.llm = AzureChatOpenAI(
        #     temperature=0,
        #     deployment_name=os.environ["AZURE_OPENAI_ENGINE"],
        #     openai_api_base=os.environ["AZURE_OPENAI_ENDPOINT"],
        #     openai_api_key=os.environ["AZURE_OPENAI_KEY"],
        #     openai_api_version=os.environ["OPENAI_API_VERSION"],
        #     openai_api_type=os.environ['OPENAI_API_TYPE'],
        # )

    # ------------------------------------------------------------------
    # Top-level retrieval entry points
    # ------------------------------------------------------------------

    async def getTopContexts(
        self,
        search_query: str,
        course_name: str,
        doc_groups: List[str] | None = None,
        top_n: int = 100,
        conversation_id: str = "",
    ) -> Union[List[Dict], str]:
        """Here's a summary of the work.

          /GET arguments
          course name (optional) str: A json response with TBD fields.

        Returns
        JSON: A json response with TBD fields. See main.py:getTopContexts docs.
        or
        String: An error message with traceback.
        """
        if doc_groups is None:
            doc_groups = []
        try:
            start_time_overall = time.monotonic()

            # Resolve per-project connections and embedding client.
            # Both Qdrant-backed and pgvector cases return a VectorDatabase;
            # embedding config now lives on its own column so it works for
            # pgvector-only projects too.
            project_sql_db = self.conn_manager.get_documents_sql_db(course_name)
            try:
                embedding_client, query_instruction = self._resolve_embedding_client(
                    course_name
                )
            except ValueError as e:
                # Misconfigured embedding_config (e.g. provider=ollama with
                # no base_url). Surface a clean user-actionable message
                # instead of a 500 with a raw traceback.
                msg = f"Embedding configuration error for course '{course_name}': {e}"
                print(msg)
                self.sentry.capture_exception(e)
                return msg

            # Create tasks for parallel execution
            with self.thread_pool_executor as executor:
                loop = asyncio.get_event_loop()
                tasks = [
                    loop.run_in_executor(
                        executor, project_sql_db.getDisabledDocGroups, course_name
                    ),
                    loop.run_in_executor(
                        executor, project_sql_db.getPublicDocGroups, course_name
                    ),
                    loop.run_in_executor(
                        executor,
                        self._embed_query_and_measure_latency,
                        search_query,
                        embedding_client,
                        query_instruction,
                    ),
                ]

            (
                disabled_doc_groups_response,
                public_doc_groups_response,
                user_query_embedding,
            ) = await asyncio.gather(*tasks)

            disabled_doc_groups = [
                doc_group for doc_group in disabled_doc_groups_response["data"]
            ]
            public_doc_groups = public_doc_groups_response["data"]

            time_for_parallel_operations = time.monotonic() - start_time_overall
            start_time_vector_search = time.monotonic()

            # Perform vector search with conversation filter
            found_docs: list[Document] = self.vector_search(
                search_query=search_query,
                course_name=course_name,
                doc_groups=doc_groups,
                user_query_embedding=user_query_embedding,
                disabled_doc_groups=disabled_doc_groups,
                public_doc_groups=public_doc_groups,
                top_n=top_n,
                conversation_id=conversation_id,
            )

            time_to_retrieve_docs = time.monotonic() - start_time_vector_search

            valid_docs = []
            for doc in found_docs:
                valid_docs.append(doc)

            print(
                f"Course: {course_name} ||| search_query: {search_query}\n"
                f"⏰ Runtime of getTopContexts: {(time.monotonic() - start_time_overall):.2f} seconds\n"
                f"Runtime for parallel operations: {time_for_parallel_operations:.2f} seconds, "
                f"Runtime to complete vector_search: {time_to_retrieve_docs:.2f} seconds"
            )
            if len(valid_docs) == 0:
                return []

            self.posthog.capture(
                event_name="getTopContexts_success_DI",
                properties={
                    "user_query": search_query,
                    "course_name": course_name,
                    # "total_tokens_used": token_counter,
                    "total_contexts_used": len(valid_docs),
                    "total_unique_docs_retrieved": len(found_docs),
                    "getTopContext_total_latency_sec": time.monotonic()
                    - start_time_overall,
                },
            )

            return self.format_for_json(valid_docs)
        except Exception as e:
            # return full traceback to front end
            err: str = f"ERROR: In /getTopContexts. Course: {course_name} ||| search_query: {search_query}\nTraceback: {traceback.print_exc} \n{e}"  # type: ignore
            traceback.print_exc()
            print(err)
            self.sentry.capture_exception(e)
            return err

    def getAll(
        self,
        course_name: str,
    ):
        """Get all course materials based on course name.
        Args:
            course_name (as uploaded on database)
        Returns:
            list of dictionaries with distinct s3 path, readable_filename and course_name, url, base_url.
        """

        project_sql_db = self.conn_manager.get_documents_sql_db(course_name)
        response = project_sql_db.getAllMaterialsForCourse(course_name)

        data = response["data"]
        unique_combinations = set()
        distinct_dicts = []

        for item in data:
            combination = (
                item["s3_path"],
                item["readable_filename"],
                item["course_name"],
                item["url"],
                item["base_url"],
            )
            if combination not in unique_combinations:
                unique_combinations.add(combination)
                distinct_dicts.append(item)

        return distinct_dicts

    # ------------------------------------------------------------------
    # LLM monitor / safety
    # ------------------------------------------------------------------

    def llm_monitor_message(
        self, course_name: str, conversation_id: str, user_email: str, model_name: str
    ) -> str:
        """
        Will store categories in DB, send email if an alert is triggered.
        """
        import json

        from ollama import Client as OllamaClient

        from ai_ta_backend.utils.email.send_transactional_email import send_email

        client = OllamaClient(
            os.environ["OLLAMA_SERVER_URL"], api_key=os.environ["NCSA_HOSTED_API_KEY"]
        )

        # Messages live on the host main DB, not on a project's external DB.
        project_sql_db = self.conn_manager.get_sql_db()
        response = project_sql_db.getMessagesFromConvoID(conversation_id)
        messages = response["data"] if isinstance(response, dict) else response.data

        llm_monitor_model = "llama-guard3:8b"

        # Map the preset categories
        safety_categories = {
            "S1": "Violent Crimes",
            "S2": "Non-Violent Crimes",
            "S3": "Sex-Related Crimes",
            "S4": "Child Sexual Exploitation",
            "S5": "Defamation",
            "S6": "Specialized Advice",
            "S7": "Privacy",
            "S8": "Intellectual Property",
            "S9": "Indiscriminate Weapons",
            "S10": "Hate",
            "S11": "Suicide & Self-Harm",
            "S12": "Sexual Content",
            "S13": "Elections",
        }

        # Categories to exclude from triggering alerts (Privacy, Defamation, and Specialized Advice, and Elections)
        excluded_categories = {"S5", "S6", "S7", "S13"}

        # Analyze each message using LLM
        for message in messages:
            if message.get("llm-monitor-tags"):
                continue

            message_content = message["content_text"]

            if message.get("role"):
                message_content = (
                    "Message from " + message.get("role") + ":\n" + message_content
                )

            analysis_result = client.chat(
                model=llm_monitor_model,
                messages=[{"role": "user", "content": message_content}],
            )

            response_content = analysis_result["message"]["content"]

            # Prepare default LLM monitor tags
            llm_monitor_tags = {"llm_monitor_model": llm_monitor_model}

            # Identify triggered categories
            triggered_categories = []
            for category_code, category_name in safety_categories.items():
                if category_code in response_content:
                    triggered_categories.append(category_name)

            # Analyze if the message should be considered unsafe for alerting
            alert_categories = [
                cat
                for cat, code in zip(
                    triggered_categories,
                    [
                        code
                        for code in safety_categories.keys()
                        if safety_categories[code] in triggered_categories
                    ],
                )
                if code not in excluded_categories
            ]

            # Assign tags to unsafe messages and send email when necessary
            if "unsafe" in response_content.lower() and alert_categories:
                llm_monitor_tags["status"] = "unsafe"
                llm_monitor_tags["triggered_categories"] = ", ".join(
                    triggered_categories
                )

                # Prepare alert email only if there are non-excluded categories
                if alert_categories:
                    alert_body = "\n".join(
                        [
                            "LLM Monitor Alert",
                            "------------------------",
                            f"Course Name: {course_name}",
                            f"User Email: {user_email}",
                            f"Conversation Model Name: {model_name}",
                            f"LLM Monitor Model Name: {llm_monitor_model}",
                            f"Convo ID: {conversation_id}",
                            "------------------------",
                            f"Responsible Role: {message.get('role')}",
                            f"Categories: {', '.join(alert_categories)}",
                            "------------------------",
                            f"Message Content:\n{json.dumps(message_content, indent=2)}",
                            "",
                        ]
                    )

                    message_id = message.get("id")
                    print(f"LLM Monitor Alert Triggered! Message ID: {message_id}")

                    send_email(
                        subject=f"LLM Monitor Alert - {', '.join(alert_categories)}",
                        body_text=alert_body,
                        sender="hi@uiuc.chat",
                        recipients=["hbroome@illinois.edu", "rohan13@illinois.edu"],
                        bcc_recipients=[],
                    )
            else:
                llm_monitor_tags["status"] = "safe"

            # Update llm_monitor_tags in messages database
            message_id = message.get("id")
            project_sql_db.updateMessageFromLlmMonitor(message_id, llm_monitor_tags)

        return "Success"

    # ------------------------------------------------------------------
    # Deletion
    # ------------------------------------------------------------------

    def delete_data(self, course_name: str, s3_path: str, source_url: str):
        """Delete file from S3, the vector store (Qdrant or pgvector), and Supabase."""
        print(f"Deleting data for course {course_name}")
        # add delete from doc map logic here
        try:
            # Resolve per-project S3 client/bucket.
            s3, bucket_name = self.conn_manager.get_s3_client(course_name)

            if bucket_name is None:
                raise ValueError("S3 bucket environment variable is not set")

            identifier_key, identifier_value = (
                ("s3_path", s3_path) if s3_path else ("url", source_url)
            )
            logging.info(
                f"Deleting {identifier_value} from S3, vector store, and Database using {identifier_key}"
            )

            # Delete from S3
            if identifier_key == "s3_path":
                self.delete_from_s3(s3, bucket_name, s3_path)
            self.delete_from_qdrant(identifier_key, identifier_value, course_name)
            self.delete_from_nomic_and_database(
                course_name, identifier_key, identifier_value
            )
            return "Success"

        except Exception as e:
            err: str = f"ERROR IN delete_data: Traceback: {traceback.extract_tb(e.__traceback__)}❌❌ Error in {inspect.currentframe().f_code.co_name}:{e}"  # type: ignore
            print(err)
            self.sentry.capture_exception(e)
            return err

    def delete_from_s3(self, s3, bucket_name: str, s3_path: str):
        try:
            print("Deleting from S3")
            response = s3.delete_object(bucket_name, s3_path)
            print(f"AWS response: {response}")
        except Exception as e:
            print("Error in deleting file from s3:", e)
            self.sentry.capture_exception(e)

    def delete_from_qdrant(
        self, identifier_key: str, identifier_value: str, course_name: str
    ):
        """Delete from the project's vector store (Qdrant or pgvector).

        Routes through ``ConnectionManager.get_vector_db`` so per-project
        Qdrant overrides hit the right cluster, and pgvector-default projects
        delete from the host embeddings table.
        """
        try:
            print("Deleting from vector store (Qdrant or pgvector)")
            vdb = self.conn_manager.get_vector_db(course_name)
            engine_kind = self.conn_manager.get_vector_engine_kind(course_name)

            if engine_kind == "qdrant":
                collection_name = vdb.qdrant_config.get(
                    "default_collection", os.environ.get("QDRANT_COLLECTION_NAME", "")
                )
            else:
                # pgvector ignores collection_name; pass a sentinel for logs only.
                collection_name = os.environ.get("QDRANT_COLLECTION_NAME", "")

            response = vdb.delete_data(
                collection_name, identifier_key, identifier_value
            )
            print(f"Vector store response: {response}")
        except Exception as e:
            if "timed out" in str(e):
                # Timed out is fine. Still deletes.
                pass
            else:
                print("Error in deleting file from vector store:", e)
                self.sentry.capture_exception(e)

    def getTopContextsWithMQR(
        self, search_query: str, course_name: str, token_limit: int = 4_000
    ) -> Union[List[Dict], str]:
        """
        New info-retrieval pipeline that uses multi-query retrieval + filtering + reciprocal rank fusion + context padding.
        1. Generate multiple queries based on the input search query.
        2. Retrieve relevant docs for each query.
        3. Filter the relevant docs based on the user query and pass them to the rank fusion step.
        4. [CANCELED BEC POINTLESS] Rank the docs based on the relevance score.
        5. Parent-doc-retrieval: Pad just the top 5 docs with expanded context from the original document.
        """
        raise NotImplementedError(
            "Method deprecated for performance reasons. Hope to bring back soon."
        )

    def delete_from_nomic_and_database(
        self, course_name: str, identifier_key: str, identifier_value: str
    ):
        try:
            print(
                f"Database Delete. course: {course_name} using {identifier_key}: {identifier_value}"
            )
            project_sql_db = self.conn_manager.get_documents_sql_db(course_name)
            response = project_sql_db.deleteMaterialsForCourseAndKeyAndValue(
                course_name, identifier_key, identifier_value
            )
        except Exception as e:
            print(f"Database Error in delete. {identifier_key}: {identifier_value}", e)
            self.sentry.capture_exception(e)

    # ------------------------------------------------------------------
    # Vector search
    # ------------------------------------------------------------------

    def vector_search(
        self,
        search_query,
        course_name,
        doc_groups: List[str],
        user_query_embedding,
        disabled_doc_groups,
        public_doc_groups,
        top_n: int = 100,
        conversation_id: str = "",
    ):
        """
        Search the vector database for a given query, course name, and document groups.
        Includes optional conversation-specific filtering.

        Engine routing (resolved per call by ConnectionManager):
          * Project with active ``qdrant_config``  -> external Qdrant
          * Otherwise                              -> pgvector
            (host pgvector by default; per-project external pg when the
            project has ``database_config``)

        The Qdrant path uses ``VectorDatabase.execute_search`` which is
        config-driven (single- vs multi-collection) — the resolver-supplied
        ``qdrant_config`` is what makes vyriad/cropwizard/pubmed/patents-style
        multi-collection setups work without hardcoded branches. The pgvector
        path goes through ``PgVectorStore.search`` with the same filter
        semantics as the frontend ``vectorSearchWithDrizzle``.
        """
        if doc_groups is None:
            doc_groups = []

        if disabled_doc_groups is None:
            disabled_doc_groups = []

        if public_doc_groups is None:
            public_doc_groups = []

        # Capture the search invoked event to PostHog
        self._capture_search_invoked_event(search_query, course_name, doc_groups)

        # Perform the vector search
        start_time_vector_search = time.monotonic()

        # Dynamic resolution: use ConnectionManager to get the right VectorDatabase
        vdb = self.conn_manager.get_vector_db(course_name)
        engine_kind = self.conn_manager.get_vector_engine_kind(course_name)

        if engine_kind == "qdrant":
            search_filter = vdb._create_search_filter(
                course_name, doc_groups, disabled_doc_groups, public_doc_groups
            )
            if conversation_id:
                chat_filter = vdb._create_conversation_search_filter(conversation_id)
                combined_filter = models.Filter(should=[search_filter, chat_filter])
                search_results = vdb.execute_search(
                    combined_filter, user_query_embedding, course_name, top_n
                )
            else:
                search_results = vdb.execute_search(
                    search_filter, user_query_embedding, course_name, top_n
                )
        else:
            # pgvector: filter shape is built inside PgVectorStore.search.
            search_results = vdb.execute_search(
                None,
                user_query_embedding,
                course_name,
                top_n,
                pgvector_args={
                    "doc_groups": doc_groups,
                    "disabled_doc_groups": disabled_doc_groups,
                    "public_doc_groups": public_doc_groups,
                    "conversation_id": conversation_id or None,
                },
            )

        self.qdrant_latency_sec = time.monotonic() - start_time_vector_search

        # Process the search results by extracting the page content and metadata
        start_time_process_search_results = time.monotonic()
        found_docs = self._process_search_results(search_results, course_name)
        time_for_process_search_results = (
            time.monotonic() - start_time_process_search_results
        )

        # Capture the search succeeded event to PostHog with the vector scores
        start_time_capture_search_succeeded_event = time.monotonic()
        self._capture_search_succeeded_event(search_query, course_name, search_results)
        time_for_capture_search_succeeded_event = (
            time.monotonic() - start_time_capture_search_succeeded_event
        )

        print(
            f"Runtime for embedding query: {self.openai_embedding_latency:.2f} seconds\n"
            f"Runtime for vector search: {self.qdrant_latency_sec:.2f} seconds\n"
            f"Runtime for process search results: {time_for_process_search_results:.2f} seconds\n"
            f"Runtime for capture search succeeded event: {time_for_capture_search_succeeded_event:.2f} seconds"
        )
        return found_docs

    def _embed_query_and_measure_latency(
        self, search_query, embedding_client, query_instruction: str | None = None
    ):
        openai_start_time = time.monotonic()
        text_to_embed = search_query

        # If using a Qwen embedding model via OpenAI-compatible API, prefix the instruction for queries only.
        try:
            model_name = getattr(embedding_client, "model", self.embedding_model)
        except Exception:
            model_name = self.embedding_model

        if (
            query_instruction
            and isinstance(embedding_client, OpenAIEmbeddings)
            and "qwen" in str(model_name).lower()
        ):
            text_to_embed = f"Instruct: {query_instruction}\nQuery:{search_query}"

        user_query_embedding = embedding_client.embed_query(text_to_embed)
        self.openai_embedding_latency = time.monotonic() - openai_start_time
        return user_query_embedding

    def _embed_document(self, text: str, embedding_client) -> list[float]:
        """Embed a document chunk without applying any query-instruction prefix.

        Why: Qwen embedding models require the ``Instruct: …\\nQuery:`` prefix
        for queries only. Documents must be embedded raw, otherwise retrieval
        quality degrades silently.
        """
        return embedding_client.embed_documents([text])[0]

    def _resolve_embedding_client(self, project_name: str):
        """Resolve the embedding client and query instruction for a project.

        Reads the top-level ``embedding_config`` column on
        ``project_external_connections`` (decrypted via ConnectionManager).
        Backward-compat: also honors a legacy ``embedding`` key nested inside
        ``qdrant_config`` so existing Qdrant projects keep working until they
        migrate to the new column.

        Falls back to the default (env-based) embedding client when no
        per-project embedding config is present.

        Returns:
            (embedding_client, query_instruction)
        """
        embedding_cfg = self.conn_manager._get_decrypted_field(
            project_name, "embedding_config"
        )
        if not embedding_cfg:
            qdrant_cfg = self.conn_manager._get_decrypted_field(
                project_name, "qdrant_config"
            )
            embedding_cfg = (qdrant_cfg or {}).get("embedding") if qdrant_cfg else None

        if not embedding_cfg:
            return self.embeddings, self.qwen_query_instruction

        provider = embedding_cfg.get("provider", "openai")
        if provider not in ALLOWED_EMBEDDING_PROVIDERS:
            raise ValueError(
                f"Unsupported embedding provider {provider!r}. "
                f"Allowed: {ALLOWED_EMBEDDING_PROVIDERS}"
            )
        model = embedding_cfg.get("model", self.embedding_model)
        query_instruction = embedding_cfg.get(
            "query_instruction", self.qwen_query_instruction
        )

        if provider == "ollama":
            base_url = embedding_cfg.get("base_url") or os.environ.get(
                "OLLAMA_SERVER_URL"
            )
            if not base_url:
                raise ValueError(
                    "embedding_config provider='ollama' requires base_url "
                    "(or set OLLAMA_BASE_URL / OLLAMA_SERVER_URL)"
                )
            client = OllamaEmbeddings(
                base_url=base_url,
                model=model,
            )
        else:
            api_key = embedding_cfg.get("api_key", self.openai_api_key)
            api_base = embedding_cfg.get("api_base", self.openai_api_base)
            kwargs = dict(
                model=model,
                openai_api_key=api_key,
                openai_api_base=api_base,
            )
            if model != "text-embedding-ada-002":
                kwargs["tiktoken_enabled"] = False
            client = OpenAIEmbeddings(**kwargs)

        return client, query_instruction

    def _capture_search_invoked_event(self, search_query, course_name, doc_groups):
        self.posthog.capture(
            event_name="vector_search_invoked",
            properties={
                "user_query": search_query,
                "course_name": course_name,
                "doc_groups": doc_groups,
            },
        )

    def _process_search_results(self, search_results, course_name):
        found_docs: list[Document] = []
        for d in search_results:
            try:
                metadata = d.payload
                # print(f"Metadata: {metadata}")
                page_content = metadata["page_content"]
                del metadata["page_content"]
                if (
                    "pagenumber" not in metadata.keys()
                    and "pagenumber_or_timestamp" in metadata.keys()
                ):
                    metadata["pagenumber"] = metadata["pagenumber_or_timestamp"]

                found_docs.append(
                    Document(page_content=page_content, metadata=metadata)
                )
            except Exception as e:
                print(
                    f"Error in vector_search(), for course: `{course_name}`. Error: {e}"
                )
                self.sentry.capture_exception(e)
        return found_docs

    def _capture_search_succeeded_event(
        self, search_query, course_name, search_results
    ):
        vector_score_calc_latency_sec = time.monotonic()
        self.posthog.capture(
            event_name="vector_search_succeeded",
            properties={
                "user_query": search_query,
                "course_name": course_name,
                "qdrant_latency_sec": self.qdrant_latency_sec,
                "openai_embedding_latency_sec": self.openai_embedding_latency,
                "vector_score_calculation_latency_sec": time.monotonic()
                - vector_score_calc_latency_sec,
            },
        )

    def _calculate_vector_scores(self, search_results):
        max_vector_score = 0
        min_vector_score = 0
        total_vector_score = 0
        for result in search_results:
            max_vector_score = max(max_vector_score, result.score)
            min_vector_score = min(min_vector_score, result.score)
            total_vector_score += result.score
        avg_vector_score = (
            total_vector_score / len(search_results) if search_results else 0
        )
        return max_vector_score, min_vector_score, avg_vector_score

    def format_for_json(self, found_docs: List[Document]) -> List[Dict]:
        """Format documents into JSON-serializable dictionaries.

        Args:
            found_docs: List of Document objects containing page content and metadata

        Returns:
            List of dictionaries with text content and metadata fields
        """
        return [
            {
                "text": doc.page_content,
                "readable_filename": doc.metadata["readable_filename"],
                "course_name ": doc.metadata["course_name"],
                # OPTIONAL
                "s3_path": doc.metadata.get("s3_path"),
                "pagenumber": doc.metadata.get(
                    "pagenumber"
                ),  # Handles both old and new schema
                "url": doc.metadata.get("url"),
                "base_url": doc.metadata.get("base_url"),
                "doc_groups": doc.metadata.get("doc_groups"),
            }
            for doc in found_docs
        ]

    # ------------------------------------------------------------------
    # Conversation stats & analytics (host main DB)
    # ------------------------------------------------------------------

    def getConversationStats(
        self, course_name: str, from_date: str = "", to_date: str = ""
    ):
        """
        Fetches conversation data from the database and groups them by day, hour, and weekday.

        Args:
            course_name (str): Name of the course
            from_date (str, optional): Start date in ISO format
            to_date (str, optional): End date in ISO format
        """
        try:
            # Conversations live on the host main DB, not on a project's external DB.
            project_sql_db = self.conn_manager.get_sql_db()
            conversations, total_count = (
                project_sql_db.getConversationsCreatedAtByCourse(
                    course_name, from_date, to_date
                )
            )

            response_data = {
                "per_day": {},
                "per_hour": {str(hour): 0 for hour in range(24)},
                "per_weekday": {
                    day: 0
                    for day in [
                        "Monday",
                        "Tuesday",
                        "Wednesday",
                        "Thursday",
                        "Friday",
                        "Saturday",
                        "Sunday",
                    ]
                },
                "heatmap": {
                    day: {str(hour): 0 for hour in range(24)}
                    for day in [
                        "Monday",
                        "Tuesday",
                        "Wednesday",
                        "Thursday",
                        "Friday",
                        "Saturday",
                        "Sunday",
                    ]
                },
                "total_count": 0,
            }

            if not conversations:
                return response_data

            central_tz = pytz.timezone("America/Chicago")
            grouped_data = {
                "per_day": defaultdict(int),
                "per_hour": defaultdict(int),
                "per_weekday": defaultdict(int),
                "heatmap": defaultdict(lambda: defaultdict(int)),
            }

            for created_at in conversations:
                try:
                    parsed_date = created_at.astimezone(central_tz)

                    day = parsed_date.date()
                    hour = parsed_date.hour
                    day_of_week = parsed_date.strftime("%A")

                    grouped_data["per_day"][str(day)] += 1
                    grouped_data["per_hour"][str(hour)] += 1
                    grouped_data["per_weekday"][day_of_week] += 1
                    grouped_data["heatmap"][day_of_week][str(hour)] += 1
                except Exception as e:
                    print(f"Error processing record: {str(e)}")
                    continue

            return {
                "per_day": dict(grouped_data["per_day"]),
                "per_hour": {str(k): v for k, v in grouped_data["per_hour"].items()},
                "per_weekday": dict(grouped_data["per_weekday"]),
                "heatmap": {
                    day: {str(h): count for h, count in hours.items()}
                    for day, hours in grouped_data["heatmap"].items()
                },
                "total_count": total_count,
            }

        except Exception as e:
            print(f"Error in getConversationStats for course {course_name}: {str(e)}")
            self.sentry.capture_exception(e)
            return {
                "per_day": {},
                "per_hour": {str(hour): 0 for hour in range(24)},
                "per_weekday": {
                    day: 0
                    for day in [
                        "Monday",
                        "Tuesday",
                        "Wednesday",
                        "Thursday",
                        "Friday",
                        "Saturday",
                        "Sunday",
                    ]
                },
                "heatmap": {
                    day: {str(hour): 0 for hour in range(24)}
                    for day in [
                        "Monday",
                        "Tuesday",
                        "Wednesday",
                        "Thursday",
                        "Friday",
                        "Saturday",
                        "Sunday",
                    ]
                },
                "total_count": 0,
            }

    def getProjectStats(self, project_name: str) -> ProjectStats:
        """
        Get statistics for a project.

        Args:
            project_name (str)

        Returns:
            ProjectStats: TypedDict containing:
                - total_messages (int): Total number of messages
                - total_conversations (int): Total number of conversations
                - unique_users (int): Number of unique users
                - avg_conversations_per_user (float): Average conversations per user
                - avg_messages_per_user (float): Average messages per user
                - avg_messages_per_conversation (float): Average messages per conversation
        """
        # Project stats live on the host main DB, not on a project's external DB.
        return self.conn_manager.get_sql_db().getProjectStats(project_name)

    def getWeeklyTrends(self, project_name: str) -> List[WeeklyMetric]:
        """
        Get weekly trends for a project, showing percentage changes in metrics.

        Args:
            project_name (str): Name of the project

        Returns:
            List[WeeklyMetric]: List of metrics with their current week value,
            previous week value, and percentage change.
        """
        # Weekly trends are computed from host conversation/message tables.
        return self.conn_manager.get_sql_db().getWeeklyTrends(project_name)

    def getModelUsageCounts(self, project_name: str) -> List[ModelUsage]:
        """
        Get counts of different models used in conversations for a project.

        Args:
            project_name (str): Name of the project

        Returns:
            List[ModelUsage]: List of model usage statistics containing model name,
            count and percentage of total usage
        """
        try:
            # Model usage counts are computed from host conversation/message tables.
            return self.conn_manager.get_sql_db().getModelUsageCounts(project_name)

        except Exception as e:
            print(f"Error fetching model usage counts for {project_name}: {str(e)}")
            self.sentry.capture_exception(e)
            return []

    def _create_conversation_filter(self, conversation_id: str):
        """Create a Qdrant filter for conversation-specific content."""
        return models.Filter(
            must=[
                models.FieldCondition(
                    key="conversation_id",
                    match=models.MatchValue(value=conversation_id),
                )
            ]
        )

    # ------------------------------------------------------------------
    # Chat-file ingest (synchronous)
    # ------------------------------------------------------------------

    def process_chat_file_sync(
        self,
        conversation_id: str,
        s3_path: str,
        course_name: str,
        readable_filename: str,
        user_id: str = None,
        is_chat_upload: bool = True,
    ):
        """
        Synchronous chat file processor - handles all allowed file types without Beam dependencies.
        Supported types: html, py, pdf, txt, md, srt, vtt, docx, ppt, pptx, xlsx, xls, xlsm,
        xlsb, xltx, xltm, xlt, xml, xlam, xla, xlw, xlr, csv, png, jpg
        """
        try:
            # Download file from S3 to process locally (per-project bucket/client).
            with tempfile.NamedTemporaryFile(
                delete=False, suffix=Path(s3_path).suffix
            ) as tmp_file:
                s3, bucket_name = self.conn_manager.get_s3_client(course_name)
                s3.download_fileobj(bucket_name, s3_path, tmp_file)
                tmp_file.flush()

                # Determine file type and extract content
                file_extension = Path(s3_path).suffix.lower()
                mime_type = str(mimetypes.guess_type(tmp_file.name, strict=False)[0])
                mime_category = (
                    mime_type.split("/")[0] if "/" in mime_type else mime_type
                )

                # Validate file type is allowed
                allowed_extensions = [
                    ".html",
                    ".py",
                    ".pdf",
                    ".txt",
                    ".md",
                    ".srt",
                    ".vtt",
                    ".docx",
                    ".ppt",
                    ".pptx",
                    ".xlsx",
                    ".xls",
                    ".xlsm",
                    ".xlsb",
                    ".xltx",
                    ".xltm",
                    ".xlt",
                    ".xml",
                    ".xlam",
                    ".xla",
                    ".xlw",
                    ".xlr",
                    ".csv",
                    ".png",
                    ".jpg",
                ]

                if file_extension not in allowed_extensions:
                    return {
                        "success": False,
                        "chunks_created": 0,
                        "error": f"File type {file_extension} not supported. Allowed types: {', '.join(allowed_extensions)}",
                    }

                # Extract text based on file type
                text_content = self._extract_file_content(
                    tmp_file.name, file_extension, mime_category
                )

                if text_content and text_content.strip():
                    # Store in vector database with conversation_id
                    result = self._store_conversation_content(
                        text_content,
                        conversation_id,
                        course_name,
                        readable_filename,
                        s3_path,
                    )

                    return result
                else:
                    return {
                        "success": False,
                        "chunks_created": 0,
                        "error": "No content extracted from file",
                    }

        except Exception as e:
            traceback.print_exc()
            return {
                "success": False,
                "chunks_created": 0,
                "error": f"Processing error: {str(e)}",
            }

    def _extract_file_content(
        self, file_path: str, file_extension: str, mime_category: str
    ) -> str:
        """
        Extract content from allowed file types only.
        """
        try:
            # Route to appropriate extraction method based on file extension
            if file_extension == ".pdf":
                return self._extract_pdf_content(file_path)
            elif file_extension in [".docx"]:
                return self._extract_docx_content(file_path)
            elif file_extension in [".txt", ".md"]:
                return self._extract_txt_content(file_path)
            elif file_extension == ".html":
                return self._extract_html_content(file_path)
            elif file_extension == ".csv":
                return self._extract_csv_content(file_path)
            elif file_extension in [
                ".xlsx",
                ".xls",
                ".xlsm",
                ".xlsb",
                ".xltx",
                ".xltm",
                ".xlt",
                ".xml",
                ".xlam",
                ".xla",
                ".xlw",
                ".xlr",
            ]:
                return self._extract_excel_content(file_path)
            elif file_extension in [".ppt", ".pptx"]:
                return self._extract_ppt_content(file_path)
            elif file_extension in [".png", ".jpg"]:
                return self._extract_image_content(file_path)
            elif file_extension in [".srt", ".vtt"]:
                return self._extract_subtitle_content(file_path)
            elif file_extension == ".py":
                return self._extract_python_content(file_path)
            else:
                return f"Unsupported file type: {file_extension}"

        except Exception as e:
            return f"Error extracting content: {e}"

    def _extract_pdf_content(self, file_path: str) -> str:
        """Extract text from PDF using available libraries."""
        try:
            # Try PyMuPDF first (most reliable)
            try:
                import fitz

                doc = fitz.open(file_path)
                text_content = []

                for i, page in enumerate(doc):
                    text = (
                        page.get_text().encode("utf8").decode("utf8", errors="ignore")
                    )
                    if text.strip():
                        text_content.append(text)

                doc.close()

                full_text = "\n".join(text_content)

                if full_text.strip():
                    return full_text
                else:
                    pass  # No text found, continue to fallback

            except ImportError:
                pass  # PyMuPDF not available, continue to fallback
            except Exception:
                pass

            # Try pdfplumber as fallback
            try:
                import pdfplumber

                text_content = []
                with pdfplumber.open(file_path) as pdf:
                    for page_num, page in enumerate(pdf.pages, 1):
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            text_content.append(page_text)

                full_text = "\n\n".join(text_content)

                if full_text.strip():
                    return full_text
                else:
                    return self._ocr_pdf_content(file_path)

            except ImportError:
                return self._ocr_pdf_content(file_path)
            except Exception:
                return self._ocr_pdf_content(file_path)

        except Exception as e:
            return f"PDF processing error: {e}"

    def _extract_docx_content(self, file_path: str) -> str:
        """Extract text from DOCX file."""
        try:
            # Try python-docx first
            try:
                from docx import Document

                doc = Document(file_path)
                text_content = []

                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        text_content.append(paragraph.text)

                # Extract from tables
                for table in doc.tables:
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            text_content.append(" | ".join(row_text))

                full_text = "\n".join(text_content)

                if full_text.strip():
                    return full_text
                else:
                    pass

            except ImportError:
                pass
            except Exception:
                pass

            # Try langchain as fallback
            try:
                from langchain.document_loaders import Docx2txtLoader

                loader = Docx2txtLoader(file_path)
                documents = loader.load()
                text_content = "\n".join([doc.page_content for doc in documents])

                if text_content.strip():
                    return text_content
                else:
                    return "No content found in DOCX file"

            except ImportError:
                return "DOCX processing libraries not available"
            except Exception as e:
                return f"DOCX processing error: {e}"

        except Exception as e:
            return f"DOCX processing error: {e}"

    def _store_conversation_content(
        self,
        text_content: str,
        conversation_id: str,
        course_name: str,
        readable_filename: str,
        s3_path: str,
    ) -> dict:
        """Store extracted content in vector database with conversation_id.

        Engine-aware: writes go through ``VectorDatabase.upsert_main_collection``
        which dispatches to Qdrant or pgvector based on the resolver.
        """
        try:
            # Split text into chunks for better retrieval
            from langchain.text_splitter import RecursiveCharacterTextSplitter

            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len,
            )

            chunks = text_splitter.split_text(text_content)

            # Resolve per-project vector DB and embedding client
            vdb = self.conn_manager.get_vector_db(course_name)
            embedding_client, _ = self._resolve_embedding_client(course_name)

            ids: List[str] = []
            vectors: List[List[float]] = []
            payloads: List[dict] = []

            for i, chunk in enumerate(chunks):
                try:
                    import uuid

                    embedding = self._embed_document(chunk, embedding_client)
                    payload = {
                        "course_name": course_name,
                        "conversation_id": conversation_id,
                        "readable_filename": readable_filename,
                        "s3_path": s3_path,
                        "page_content": chunk,
                        "chunk_index": i,
                        "pagenumber": "",
                        "url": "",
                        "base_url": "",
                    }
                    ids.append(str(uuid.uuid4()))
                    vectors.append(embedding)
                    payloads.append(payload)

                except Exception as e:
                    print("Error in _store_conversation_content: ", e)
                    continue

            if ids:
                # Engine-aware ingest. ``upsert_main_collection`` routes to
                # pgvector or Qdrant based on the resolved VectorDatabase.
                vdb.upsert_main_collection(
                    ids=ids, vectors=vectors, payloads=payloads, wait=True
                )

                return {
                    "success": True,
                    "chunks_created": len(ids),
                    "total_chunks_attempted": len(chunks),
                }
            else:
                return {
                    "success": False,
                    "chunks_created": 0,
                    "total_chunks_attempted": len(chunks),
                }

        except Exception as e:
            traceback.print_exc()
            self.sentry.capture_exception(e)
            return {
                "success": False,
                "chunks_created": 0,
                "total_chunks_attempted": 0,
                "error": str(e),
            }

    def _extract_excel_content(self, file_path: str) -> str:
        """Extract text from Excel files (.xlsx, .xls, etc.) with calculated values."""
        try:
            # Try enhanced openpyxl first for better value extraction
            try:
                import openpyxl

                # Load with data_only=True to get calculated values instead of formulas
                workbook = openpyxl.load_workbook(file_path, data_only=True)
                all_text = []

                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    sheet_text = []

                    # Get the used range to avoid empty cells
                    if sheet.max_row > 0 and sheet.max_column > 0:
                        for row_num in range(1, sheet.max_row + 1):
                            row_data = []
                            for col_num in range(1, sheet.max_column + 1):
                                cell = sheet.cell(row=row_num, column=col_num)
                                cell_value = cell.value

                                if cell_value is not None:
                                    # Handle different data types
                                    if isinstance(cell_value, (int, float)):
                                        # For numeric values, include cell reference
                                        col_letter = openpyxl.utils.get_column_letter(
                                            col_num
                                        )
                                        cell_ref = f"{col_letter}{row_num}"
                                        row_data.append(
                                            f"{cell_value} (cell {cell_ref})"
                                        )
                                    elif (
                                        isinstance(cell_value, str)
                                        and cell_value.strip()
                                    ):
                                        row_data.append(cell_value.strip())
                                    else:
                                        row_data.append(str(cell_value))

                            if row_data:
                                sheet_text.append(
                                    f"Row {row_num}: " + " | ".join(row_data)
                                )

                    if sheet_text:
                        sheet_content = f"\n--- Sheet: {sheet_name} ---\n" + "\n".join(
                            sheet_text
                        )
                        all_text.append(sheet_content)

                content = "\n\n".join(all_text)

                if content.strip():
                    return content

            except ImportError:
                pass
            except Exception:
                pass

            # Fallback to langchain
            try:
                from langchain.document_loaders import UnstructuredExcelLoader

                loader = UnstructuredExcelLoader(file_path, mode="elements")
                documents = loader.load()
                text_content = "\n".join([doc.page_content for doc in documents])

                if text_content.strip():
                    return text_content
                else:
                    return "No content found in Excel file"

            except ImportError:
                return "Excel processing libraries not available"
            except Exception as e:
                return f"Excel processing error: {e}"

        except Exception as e:
            return f"Excel processing error: {e}"

    def _extract_ppt_content(self, file_path: str) -> str:
        """Extract comprehensive text from PowerPoint files (.ppt, .pptx)."""
        try:
            # Try enhanced python-pptx first for detailed extraction
            try:
                from pptx import Presentation
                from pptx.enum.shapes import MSO_SHAPE_TYPE

                prs = Presentation(file_path)
                all_content = []

                # Extract presentation title if available
                if hasattr(prs.core_properties, "title") and prs.core_properties.title:
                    all_content.append(
                        f"PRESENTATION TITLE: {prs.core_properties.title}"
                    )

                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_content = []
                    slide_content.append(f"\n=== SLIDE {slide_num} ===")

                    # Extract slide title if available
                    if slide.shapes.title and slide.shapes.title.text.strip():
                        slide_content.append(
                            f"SLIDE TITLE: {slide.shapes.title.text.strip()}"
                        )

                    # Process all shapes on the slide
                    for shape in slide.shapes:
                        try:
                            # Text content from text boxes and placeholders
                            if hasattr(shape, "text") and shape.text.strip():
                                # Skip title (already extracted)
                                if shape != slide.shapes.title:
                                    slide_content.append(
                                        f"TEXT CONTENT: {shape.text.strip()}"
                                    )

                            # Table content
                            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                                table_content = []
                                table_content.append("TABLE DATA:")
                                for row_idx, row in enumerate(shape.table.rows):
                                    row_data = []
                                    for cell in row.cells:
                                        if cell.text.strip():
                                            row_data.append(cell.text.strip())
                                    if row_data:
                                        table_content.append(
                                            f"  Row {row_idx + 1}: {' | '.join(row_data)}"
                                        )

                                if len(table_content) > 1:  # More than just header
                                    slide_content.extend(table_content)

                            # Chart content (basic info)
                            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                                slide_content.append(
                                    "CHART DETECTED: Visual data representation present"
                                )
                                if hasattr(shape, "chart") and shape.chart:
                                    try:
                                        # Try to get chart title
                                        if (
                                            hasattr(shape.chart, "chart_title")
                                            and shape.chart.chart_title
                                        ):
                                            slide_content.append(
                                                f"CHART TITLE: {shape.chart.chart_title.text_frame.text}"
                                            )
                                    except Exception:
                                        pass

                            # Group shapes (nested content)
                            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                                group_content = []
                                for group_shape in shape.shapes:
                                    if (
                                        hasattr(group_shape, "text")
                                        and group_shape.text.strip()
                                    ):
                                        group_content.append(group_shape.text.strip())
                                if group_content:
                                    slide_content.append(
                                        f"GROUPED CONTENT: {' | '.join(group_content)}"
                                    )

                            # SmartArt and other shapes with text
                            elif hasattr(shape, "text_frame") and shape.text_frame:
                                for paragraph in shape.text_frame.paragraphs:
                                    if paragraph.text.strip():
                                        slide_content.append(
                                            f"CONTENT: {paragraph.text.strip()}"
                                        )

                        except Exception:
                            continue

                    # Add slide notes if available
                    if slide.notes_slide and slide.notes_slide.notes_text_frame:
                        notes_text = slide.notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            slide_content.append(f"SLIDE NOTES: {notes_text}")

                    # Only add slide if it has content
                    if len(slide_content) > 1:  # More than just the slide header
                        all_content.extend(slide_content)

                content = "\n".join(all_content)

                if content.strip():
                    print(
                        f"PowerPoint processed successfully with python-pptx - extracted {len(content)} characters"
                    )
                    return content

            except ImportError:
                pass
            except Exception:
                pass

            # Fallback to langchain
            try:
                from langchain.document_loaders import UnstructuredPowerPointLoader

                loader = UnstructuredPowerPointLoader(file_path)
                documents = loader.load()
                text_content = "\n".join([doc.page_content for doc in documents])

                if text_content.strip():
                    print(
                        f"PowerPoint processed with langchain - extracted {len(text_content)} characters"
                    )
                    return text_content
                else:
                    return "No content found in PowerPoint file"

            except ImportError:
                return "PowerPoint processing libraries not available"
            except Exception as e:
                return f"PowerPoint processing error: {e}"

        except Exception as e:
            return f"PowerPoint processing error: {e}"

    def _extract_csv_content(self, file_path: str) -> str:
        """Extract text from CSV files."""
        try:
            # Try langchain first
            try:
                from langchain.document_loaders import CSVLoader

                loader = CSVLoader(file_path)
                documents = loader.load()
                text_content = "\n".join([doc.page_content for doc in documents])

                if text_content.strip():
                    print(
                        f"CSV processed with langchain - extracted {len(text_content)} characters"
                    )
                    return text_content

            except ImportError:
                pass
            except Exception:
                pass

            # Fallback to basic csv reading
            try:
                import csv

                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    reader = csv.reader(f)
                    rows = []
                    for row_num, row in enumerate(reader, 1):
                        if row:  # Skip empty rows
                            rows.append(f"Row {row_num}: {', '.join(row)}")

                    if rows:
                        return "\n".join(rows)
                    else:
                        return "No content found in CSV file"

            except Exception as e:
                return f"CSV processing error: {e}"

        except Exception as e:
            return f"CSV processing error: {e}"

    def _extract_txt_content(self, file_path: str) -> str:
        """Extract text from plain text files (.txt, .md)."""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()

            if content.strip():
                return content
            else:
                return "No content found in text file"

        except Exception as e:
            return f"Text file processing error: {e}"

    def _extract_html_content(self, file_path: str) -> str:
        """Extract text from HTML files."""
        try:
            # Try BeautifulSoup first
            try:
                from bs4 import BeautifulSoup

                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()

                soup = BeautifulSoup(content, "html.parser")

                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()

                text = soup.get_text()

                # Clean up whitespace
                lines = (line.strip() for line in text.splitlines())
                chunks = (
                    phrase.strip() for line in lines for phrase in line.split("  ")
                )
                text = "\n".join(chunk for chunk in chunks if chunk)

                if text.strip():
                    return text

            except ImportError:
                pass
            except Exception:
                pass

            # Fallback to basic tag removal
            try:
                import re

                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()

                # Remove HTML tags
                text = re.sub("<[^<]+?>", "", content)

                if text.strip():
                    return text
                else:
                    return "No content found in HTML file"

            except Exception as e:
                return f"HTML processing error: {e}"

        except Exception as e:
            return f"HTML processing error: {e}"

    def _extract_image_content(self, file_path: str) -> str:
        """Extract text from images (.png, .jpg) using OCR."""
        try:
            try:
                import pytesseract
                from PIL import Image

                text_content = pytesseract.image_to_string(Image.open(file_path))

                if text_content.strip():
                    return text_content
                else:
                    return "No text found in image"

            except ImportError:
                return "Image OCR libraries not available (pytesseract, PIL)"
            except Exception as e:
                return f"Image OCR error: {e}"

        except Exception as e:
            return f"Image processing error: {e}"

    def _extract_subtitle_content(self, file_path: str) -> str:
        """Extract text from subtitle files (.srt, .vtt)."""
        try:
            if file_path.endswith(".srt"):
                try:
                    import pysrt

                    subs = pysrt.open(file_path)
                    text_content = " ".join([sub.text for sub in subs])

                    if text_content.strip():
                        return text_content
                    else:
                        return "No content found in SRT file"

                except ImportError:
                    return self._extract_txt_content(file_path)
                except Exception as e:
                    return f"SRT processing error: {e}"

            else:  # VTT file
                try:
                    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                        content = f.read()

                    if content.strip():
                        return content
                    else:
                        return "No content found in VTT file"

                except Exception as e:
                    return f"VTT processing error: {e}"

        except Exception as e:
            return f"Subtitle processing error: {e}"

    def _extract_python_content(self, file_path: str) -> str:
        """Extract text from Python files (.py)."""
        try:
            try:
                from langchain.document_loaders import PythonLoader

                loader = PythonLoader(file_path)
                documents = loader.load()
                text_content = "\n".join([doc.page_content for doc in documents])

                if text_content.strip():
                    return text_content

            except ImportError:
                pass
            except Exception:
                pass

            # Fallback to basic text reading
            return self._extract_txt_content(file_path)

        except Exception as e:
            return f"Python processing error: {e}"

    def _ocr_pdf_content(self, file_path: str) -> str:
        """OCR PDF content using pdfplumber and pytesseract."""
        try:
            import pdfplumber
            import pytesseract

            text_content = []

            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    try:
                        im = page.to_image()
                        text = pytesseract.image_to_string(im.original)
                        if text.strip():
                            text_content.append(text)
                    except Exception:
                        continue

            if text_content:
                full_text = "\n\n".join(text_content)
                return full_text
            else:
                return "No text found via OCR"

        except ImportError:
            return "OCR libraries not available: No module named 'pdfplumber'"
        except Exception as e:
            return f"OCR processing error: {e}"
